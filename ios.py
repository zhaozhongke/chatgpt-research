from pptx import Presentation
from pptx.util import Inches

# 创建演示文稿实例
presentation = Presentation()

# 幻灯片标题
slide_title = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "iOS软硬件架构"
subtitle.text = "硬件和软件层次概述"

# 用于添加内容的幻灯片布局
content_slide_layout = presentation.slide_layouts[1]

# 硬件架构
slide_hardware = presentation.slides.add_slide(content_slide_layout)
title = slide_hardware.shapes.title
title.text = "硬件架构"
content = slide_hardware.placeholders[1].text_frame
content.text = "处理器：A系列处理器（基于ARM架构）\nGPU（图形处理器）：苹果自家设计\n内存：RAM（随机存取存储器）\n存储：闪存（Flash storage）\n其他硬件组件：摄像头、触摸屏、陀螺仪、加速度计、指纹识别器、Face ID等"

# 软件架构 - 层次概述
slide_software_overview = presentation.slides.add_slide(content_slide_layout)
title = slide_software_overview.shapes.title
title.text = "软件架构 - 层次概述"
content = slide_software_overview.placeholders[1].text_frame
content.text = "1. Core OS层\n2. Core Services层\n3. Media层\n4. Cocoa Touch层"

# Core OS层
slide_core_os = presentation.slides.add_slide(content_slide_layout)
title = slide_core_os.shapes.title
title.text = "Core OS层"
content = slide_core_os.placeholders[1].text_frame
content.text = "XNU内核\n底层驱动程序\n内存管理、线程管理、文件系统\n与硬件的直接交互"

# Core Services层
slide_core_services = presentation.slides.add_slide(content_slide_layout)
title = slide_core_services.shapes.title
title.text = "Core Services层"
content = slide_core_services.placeholders[1].text_frame
content.text = "Core Foundation\nCore Data\nCore Motion\n基本服务和框架，如数据管理、网络连接和设备传感器访问"

# Media层
slide_media = presentation.slides.add_slide(content_slide_layout)
title = slide_media.shapes.title
title.text = "Media层"
content = slide_media.placeholders[1].text_frame
content.text = "Core Graphics\nCore Animation\nAVFoundation\n处理图像、音频和视频的框架"

# Cocoa Touch层
slide_cocoa_touch = presentation.slides.add_slide(content_slide_layout)
title = slide_cocoa_touch.shapes.title
title.text = "Cocoa Touch层"
content = slide_cocoa_touch.placeholders[1].text_frame
content.text = "UIKit\nMapKit\nGameKit\n构建图形用户界面（GUI）和与用户互动的工具和框架\n实现多点触控、手势识别、动画以及与系统服务（如通知和定位服务）交互的高级接口"

presentation.save("iOS_软硬件架构.pptx")

